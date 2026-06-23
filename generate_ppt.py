# -*- coding: utf-8 -*-
"""
Generateur PPT professionnel — TOX21 Prediction de Toxicite Moleculaire
20 slides, design avance bleu marine / blanc / bleu electrique
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Dimensions 16:9 ──────────────────────────────────────────
SW = Inches(13.33)
SH = Inches(7.5)

# ── Palette ──────────────────────────────────────────────────
NAVY    = RGBColor(0x0f, 0x1e, 0x35)
NAVY2   = RGBColor(0x06, 0x0e, 0x1c)
BLUE    = RGBColor(0x1a, 0x56, 0xdb)
BLUE2   = RGBColor(0x1e, 0x40, 0xaf)
BLUE_L  = RGBColor(0xdb, 0xea, 0xfe)
BLUE_M  = RGBColor(0x60, 0xa5, 0xfa)
TEAL    = RGBColor(0x0d, 0x94, 0x88)
GREEN   = RGBColor(0x15, 0x80, 0x3d)
GREEN_L = RGBColor(0xdc, 0xfc, 0xe7)
RED     = RGBColor(0xb9, 0x1c, 0x1c)
RED_L   = RGBColor(0xfe, 0xe2, 0xe2)
AMBER   = RGBColor(0x92, 0x40, 0x0e)
AMBER_L = RGBColor(0xfe, 0xf3, 0xc7)
GRAY_D  = RGBColor(0x1e, 0x29, 0x3b)
GRAY_M  = RGBColor(0x47, 0x55, 0x69)
GRAY_L  = RGBColor(0xe2, 0xe8, 0xf0)
GRAY_LL = RGBColor(0xf8, 0xfa, 0xfc)
WHITE   = RGBColor(0xff, 0xff, 0xff)
YELLOW  = RGBColor(0xfb, 0xbf, 0x24)


# ══════════════════════════════════════════════════════════════
# HELPERS
# ══════════════════════════════════════════════════════════════

def add_rect(slide, x, y, w, h, fill, border=None, radius=False):
    """Rectangle (ou arrondi) avec couleur de remplissage."""
    stype = 5 if radius else 1
    shp   = slide.shapes.add_shape(stype, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if border:
        shp.line.color.rgb = border
        shp.line.width     = Pt(1.2)
    else:
        shp.line.color.rgb = fill  # invisible
    return shp


def add_txt(slide, x, y, w, h, text,
            size=14, bold=False, italic=False,
            color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    """Zone de texte simple."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text          = text
    run.font.size     = Pt(size)
    run.font.bold     = bold
    run.font.italic   = italic
    run.font.color.rgb = color
    return tb


def add_para(tf, text, size=12, bold=False, italic=False,
             color=GRAY_D, align=PP_ALIGN.LEFT, space_before=0):
    """Ajouter un paragraphe dans un text_frame existant."""
    p   = tf.add_paragraph()
    p.alignment    = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text           = text
    run.font.size      = Pt(size)
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.color.rgb = color
    return p


def txt_box_multi(slide, x, y, w, h, lines):
    """
    Zone de texte multi-lignes.
    lines = [(text, size, bold, italic, color, align, space_before), ...]
    """
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for (text, size, bold, italic, color, align, sp) in lines:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.alignment    = align
        p.space_before = Pt(sp)
        run = p.add_run()
        run.text           = text
        run.font.size      = Pt(size)
        run.font.bold      = bold
        run.font.italic    = italic
        run.font.color.rgb = color
    return tb


def slide_base(prs, bg_color=GRAY_LL):
    """Cree un slide vierge avec fond de couleur."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, 13.33, 7.5, bg_color)
    return slide


def left_sidebar(slide, num="01", label="SECTION"):
    """Barre laterale gauche marine avec numero."""
    add_rect(slide, 0, 0, 2.0, 7.5, NAVY)
    add_rect(slide, 0, 0, 2.0, 0.06, BLUE)
    add_txt(slide, 0.0, 0.25, 2.0, 1.0, num,
            size=52, bold=True, color=BLUE_M, align=PP_ALIGN.CENTER)
    add_txt(slide, 0.0, 1.2, 2.0, 0.5, label,
            size=9, bold=True, color=GRAY_L, align=PP_ALIGN.CENTER)
    add_rect(slide, 0.25, 1.75, 1.5, 0.04, BLUE)


def top_bar(slide):
    """Barre superieure bleue fine."""
    add_rect(slide, 2.0, 0, 11.33, 0.08, BLUE)


def slide_title(slide, title, subtitle=None):
    """Titre de section sous la barre."""
    add_txt(slide, 2.2, 0.15, 10.8, 0.7, title,
            size=26, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
    add_rect(slide, 2.2, 0.85, 10.8, 0.04, BLUE_L)
    if subtitle:
        add_txt(slide, 2.2, 0.88, 10.8, 0.4, subtitle,
                size=11, bold=False, color=GRAY_M, align=PP_ALIGN.LEFT)


def metric_card(slide, x, y, w, h, value, label, val_color=BLUE, bg=GRAY_LL):
    """Carte metrique avec valeur + label."""
    add_rect(slide, x, y, w, h, bg, border=GRAY_L, radius=True)
    add_txt(slide, x, y+0.1, w, h*0.55, value,
            size=28, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    add_txt(slide, x, y+h*0.62, w, h*0.38, label,
            size=9.5, bold=False, color=GRAY_M, align=PP_ALIGN.CENTER)


def bullet_item(slide, x, y, w, text, size=12, color=GRAY_D):
    """Item de liste avec point bleu."""
    add_txt(slide, x, y, 0.25, 0.35, "▶", size=9, bold=True, color=BLUE)
    add_txt(slide, x+0.25, y, w-0.25, 0.4, text, size=size, color=color)


def add_table(slide, x, y, w, h, headers, rows, col_ws=None):
    """Table avec entete marine."""
    nrows = len(rows) + 1
    ncols = len(headers)
    tbl   = slide.shapes.add_table(
        nrows, ncols, Inches(x), Inches(y), Inches(w), Inches(h)
    ).table

    # Largeurs colonnes
    if col_ws:
        for i, cw in enumerate(col_ws):
            tbl.columns[i].width = Inches(cw)

    def fmt_cell(cell, text, bg, fg, sz=10, bold=False, align=PP_ALIGN.LEFT):
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        p = cell.text_frame.paragraphs[0]
        p.alignment = align
        # Padding via XML
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        for margin in ('marL','marR','marT','marB'):
            tcPr.set(margin, str(Pt(6).pt.__trunc__() * 12700 // 100))
        run = p.add_run()
        run.text           = text
        run.font.size      = Pt(sz)
        run.font.bold      = bold
        run.font.color.rgb = fg
        return cell

    # Header
    for j, h_text in enumerate(headers):
        fmt_cell(tbl.cell(0, j), h_text, NAVY, WHITE, 10, True, PP_ALIGN.CENTER)

    # Rows
    for i, row in enumerate(rows):
        bg = GRAY_LL if i % 2 == 0 else WHITE
        for j, cell_text in enumerate(row):
            fmt_cell(tbl.cell(i+1, j), str(cell_text), bg, GRAY_D, 9.5, False,
                     PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)
    return tbl


def tag_badge(slide, x, y, text, fill=BLUE, fg=WHITE, size=9):
    """Petit badge colore."""
    add_rect(slide, x, y, len(text)*0.09+0.2, 0.28, fill, radius=True)
    add_txt(slide, x+0.05, y+0.02, len(text)*0.09+0.1, 0.26,
            text, size=size, bold=True, color=fg, align=PP_ALIGN.CENTER)


def divider_line(slide, x, y, w, color=BLUE_L):
    add_rect(slide, x, y, w, 0.03, color)


# ══════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════

def slide_01_cover(prs):
    """Slide 1 — Couverture"""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Fond naval
    add_rect(slide, 0, 0, 13.33, 7.5, NAVY)
    # Degrade simule
    add_rect(slide, 0, 4.5, 13.33, 3.0, RGBColor(0x0a, 0x14, 0x24))
    # Cercle deco haut-droite
    for i, (r_size, alpha) in enumerate([(3.5, 0.05), (2.5, 0.07), (1.5, 0.1)]):
        c = add_rect(slide, 13.33-r_size*0.7, -r_size*0.4, r_size, r_size,
                     BLUE2, radius=False)
        c.fill.fore_color.rgb = RGBColor(0x1a+i*5, 0x40+i*5, 0x9f+i*5)
        c.line.color.rgb      = RGBColor(0x1a+i*5, 0x40+i*5, 0x9f+i*5)
    # Cercle logo centre
    logo = add_rect(slide, 5.9, 0.3, 1.55, 1.55, BLUE2, radius=False)
    logo.fill.fore_color.rgb = RGBColor(0x0d, 0x2a, 0x52)
    logo.line.color.rgb      = BLUE_M
    logo.line.width          = Pt(2)
    add_txt(slide, 5.9, 0.5, 1.55, 1.1, "TOX",
            size=32, bold=True, color=BLUE_M, align=PP_ALIGN.CENTER)
    add_txt(slide, 5.9, 1.1, 1.55, 0.5, "21",
            size=18, bold=True, color=BLUE_M, align=PP_ALIGN.CENTER)

    # Titre principal
    add_txt(slide, 1.5, 2.1, 10.33, 0.85,
            "PREDICTION DE TOXICITE MOLECULAIRE",
            size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_txt(slide, 1.5, 2.95, 10.33, 0.5,
            "TOX21  —  Machine Learning Multi-Label  |  12 Cibles Biologiques",
            size=14, bold=False, color=BLUE_M, align=PP_ALIGN.CENTER)

    # Ligne separatrice
    add_rect(slide, 2.0, 3.55, 9.33, 0.04, BLUE)

    # 5 cartes stats
    stats = [("12","Cibles Bio."), ("8 014","Molecules"), ("2 063","Features"),
             ("0.783","AUC-ROC"), ("200","Arbres RF")]
    for i, (val, lbl) in enumerate(stats):
        bx = 1.5 + i * 2.1
        add_rect(slide, bx, 3.75, 1.85, 1.5, RGBColor(0x0d, 0x1b, 0x35), radius=True)
        add_rect(slide, bx, 3.75, 1.85, 1.5, RGBColor(0x0d, 0x1b, 0x35), border=BLUE, radius=True)
        add_txt(slide, bx, 3.9, 1.85, 0.75, val,
                size=26, bold=True, color=BLUE_M, align=PP_ALIGN.CENTER)
        add_txt(slide, bx, 4.7, 1.85, 0.45, lbl,
                size=10, color=GRAY_M, align=PP_ALIGN.CENTER)

    # Infos universite
    add_txt(slide, 1.5, 5.55, 10.33, 0.4,
            "Universite Hassan II  —  Faculte des Sciences Ben M'Sick  —  Master Data Sciences & Big Data",
            size=10, color=GRAY_M, align=PP_ALIGN.CENTER)
    add_txt(slide, 1.5, 5.95, 10.33, 0.35,
            "Annee Universitaire 2025 — 2026  |  Module : Machine Learning",
            size=10, bold=True, color=GRAY_L, align=PP_ALIGN.CENTER)

    # Barre bleue bas
    add_rect(slide, 0, 6.9, 13.33, 0.6, BLUE)
    add_txt(slide, 0, 6.95, 13.33, 0.45,
            "Random Forest  •  RDKit  •  SHAP  •  Streamlit  •  py3Dmol",
            size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def slide_02_toc(prs):
    """Slide 2 — Table des matieres"""
    slide = slide_base(prs)
    left_sidebar(slide, "00", "SOMMAIRE")
    top_bar(slide)
    slide_title(slide, "Sommaire", "Apercu du contenu de la presentation")

    chapters = [
        ("01", "Contexte et Motivation"),
        ("02", "Probleme Multi-Label"),
        ("03", "Dataset TOX21 — 12 Cibles"),
        ("04", "Pipeline Methodologique"),
        ("05", "Feature Engineering"),
        ("06", "Architecture du Modele"),
        ("07", "Hyperparametres et Entrainement"),
        ("08", "Resultats Globaux"),
        ("09", "Performance par Cible"),
        ("10", "Explainabilite SHAP"),
        ("11", "Interface Streamlit"),
        ("12", "Conclusion & Perspectives"),
    ]
    cols = [chapters[:6], chapters[6:]]
    for ci, col in enumerate(cols):
        bx = 2.3 + ci * 5.5
        for ri, (num, title) in enumerate(col):
            by = 1.3 + ri * 0.95
            add_rect(slide, bx, by, 5.2, 0.78, WHITE, border=GRAY_L, radius=True)
            add_rect(slide, bx, by, 0.55, 0.78, BLUE, radius=False)
            add_txt(slide, bx, by+0.18, 0.55, 0.42, num,
                    size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            add_txt(slide, bx+0.65, by+0.2, 4.45, 0.42, title,
                    size=12, bold=True if ri == 0 else False, color=NAVY)
    add_txt(slide, 2.2, 7.15, 10.8, 0.3,
            "TOX21 — Prediction de Toxicite Moleculaire  |  Master Big Data S2 — 2025/2026",
            size=8, italic=True, color=GRAY_M)


def slide_03_context(prs):
    """Slide 3 — Contexte et Motivation"""
    slide = slide_base(prs)
    left_sidebar(slide, "01", "INTRO")
    top_bar(slide)
    slide_title(slide, "Contexte et Motivation",
                "Pourquoi predire la toxicite par Machine Learning ?")

    # Colonne gauche — probleme
    add_rect(slide, 2.2, 1.25, 5.2, 5.6, WHITE, border=GRAY_L, radius=True)
    add_rect(slide, 2.2, 1.25, 5.2, 0.5, NAVY, radius=False)
    add_txt(slide, 2.3, 1.3, 5.0, 0.45, "Le Probleme",
            size=13, bold=True, color=WHITE)

    lines_l = [
        ("Les tests biologiques traditionnels (in vivo / in vitro) sont :", 11, False, GRAY_D),
        ("", 6, False, GRAY_D),
        ("⏳  Lents — plusieurs semaines par molecule", 11, True, NAVY),
        ("💰  Couteux — des milliers d'euros par test", 11, True, NAVY),
        ("🐁  Ethiquement contraints — experimentation animale", 11, True, NAVY),
        ("", 6, False, GRAY_D),
        ("La solution : predire la toxicite directement depuis la structure "
         "chimique d'une molecule, sans aucun test biologique.", 11, False, GRAY_D),
    ]
    tb = slide.shapes.add_textbox(Inches(2.35), Inches(1.85), Inches(5.0), Inches(4.8))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for (txt_s, sz, bld, clr) in lines_l:
        if first: p = tf.paragraphs[0]; first = False
        else:     p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run(); run.text = txt_s
        run.font.size = Pt(sz); run.font.bold = bld; run.font.color.rgb = clr

    # Colonne droite — programme TOX21
    add_rect(slide, 7.65, 1.25, 5.45, 5.6, WHITE, border=GRAY_L, radius=True)
    add_rect(slide, 7.65, 1.25, 5.45, 0.5, BLUE, radius=False)
    add_txt(slide, 7.75, 1.3, 5.25, 0.45, "Programme TOX21",
            size=13, bold=True, color=WHITE)

    infos = [
        ("NIH + EPA + FDA", "Co-finances par 3 agences US"),
        ("8 014", "Molecules testees"),
        ("12", "Cibles biologiques"),
        ("High-Throughput", "Screening automatise"),
        ("Public", "Dataset libre d'acces"),
    ]
    for i, (val, desc) in enumerate(infos):
        by = 1.9 + i * 0.98
        add_rect(slide, 7.75, by, 5.2, 0.82, GRAY_LL, radius=True)
        add_txt(slide, 7.85, by+0.08, 2.5, 0.38, val,
                size=16, bold=True, color=BLUE)
        add_txt(slide, 7.85, by+0.46, 5.0, 0.32, desc,
                size=10, color=GRAY_M)

    add_txt(slide, 2.2, 7.15, 10.8, 0.3,
            "TOX21 — Machine Learning  |  Master Big Data S2 — 2025/2026",
            size=8, italic=True, color=GRAY_M)


def slide_04_multilabel(prs):
    """Slide 4 — Probleme Multi-Label"""
    slide = slide_base(prs)
    left_sidebar(slide, "02", "PROBLEME")
    top_bar(slide)
    slide_title(slide, "Probleme de Classification Multi-Label",
                "Predire simultanement 12 cibles a partir d'un seul SMILES")

    # Schema central
    # SMILES box
    add_rect(slide, 2.3, 1.6, 3.5, 0.9, NAVY, radius=True)
    add_txt(slide, 2.3, 1.68, 3.5, 0.75,
            "SMILES Input\nCC(=O)Oc1ccccc1C(=O)O",
            size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Arrow
    add_rect(slide, 5.9, 1.97, 1.0, 0.06, BLUE)
    add_txt(slide, 6.65, 1.78, 0.5, 0.45, "▶", size=18, bold=True, color=BLUE)

    # Model box
    add_rect(slide, 7.0, 1.3, 2.5, 1.7, BLUE, radius=True)
    add_txt(slide, 7.0, 1.4, 2.5, 0.5, "Modele ML",
            size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_txt(slide, 7.0, 1.85, 2.5, 0.85,
            "Random Forest\nMultiOutputClassifier\n12 estimateurs",
            size=9.5, color=BLUE_L, align=PP_ALIGN.CENTER)

    # Arrow
    add_rect(slide, 9.6, 1.97, 0.8, 0.06, BLUE)
    add_txt(slide, 10.15, 1.78, 0.5, 0.45, "▶", size=18, bold=True, color=BLUE)

    # Output boxes (12 labels)
    targets_short = ["NR-AR","NR-AhR","NR-ER","NR-Aromatase","NR-PPAR",
                     "NR-AR-LBD","SR-ARE","SR-ATAD5","SR-HSE","SR-MMP",
                     "SR-p53","NR-ER-LBD"]
    colors_t = [RED]*7 + [TEAL]*5
    for i, (tgt, clr) in enumerate(zip(targets_short, colors_t)):
        col = i % 2; row = i // 2
        bx  = 10.75 + col * 1.35
        by  = 1.1  + row * 0.55
        add_rect(slide, bx, by, 1.25, 0.44, clr, radius=True)
        add_txt(slide, bx, by+0.06, 1.25, 0.32, tgt,
                size=7.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Legende desequilibre
    add_rect(slide, 2.3, 3.1, 8.5, 0.06, GRAY_L)

    info_cards = [
        (BLUE,  "Classification\nBinaire", "0 = non-toxique\n1 = toxique"),
        (AMBER, "Desequilibre\nClasses",  "3 % - 15 % positifs\nselon la cible"),
        (GREEN, "Valeurs\nManquantes",    "15 % - 40 %\nselon la cible"),
        (NAVY,  "Approche\nBinary Relevance","1 classifieur RF\nindependant / cible"),
    ]
    for i, (clr, title, desc) in enumerate(info_cards):
        bx = 2.3 + i * 2.15
        add_rect(slide, bx, 3.3, 2.0, 2.6, WHITE, border=GRAY_L, radius=True)
        add_rect(slide, bx, 3.3, 2.0, 0.08, clr)
        add_txt(slide, bx+0.1, 3.45, 1.8, 0.65, title,
                size=12, bold=True, color=clr, align=PP_ALIGN.CENTER)
        add_txt(slide, bx+0.1, 4.1, 1.8, 0.8, desc,
                size=10, color=GRAY_M, align=PP_ALIGN.CENTER)

    add_txt(slide, 2.2, 7.15, 10.8, 0.3,
            "TOX21 — Machine Learning  |  Master Big Data S2 — 2025/2026",
            size=8, italic=True, color=GRAY_M)


def slide_05_dataset(prs):
    """Slide 5 — Dataset TOX21"""
    slide = slide_base(prs)
    left_sidebar(slide, "03", "DATASET")
    top_bar(slide)
    slide_title(slide, "Dataset TOX21 — Apercu", "NIH / EPA — Benchmark officiel de toxicologie")

    # Metriques
    mets = [("8 014","Molecules"), ("12","Cibles Bio."), ("2 063","Features"),
            ("40 %","Val. manquantes max"), ("15 %","Positifs max")]
    for i, (v, l) in enumerate(mets):
        bx = 2.3 + i * 2.22
        add_rect(slide, bx, 1.15, 2.0, 1.25, WHITE, border=GRAY_L, radius=True)
        add_rect(slide, bx, 1.15, 2.0, 0.08, BLUE)
        add_txt(slide, bx, 1.3, 2.0, 0.65, v,
                size=26, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
        add_txt(slide, bx, 1.95, 2.0, 0.38, l,
                size=9, color=GRAY_M, align=PP_ALIGN.CENTER)

    # Table caracteristiques
    add_table(slide, 2.2, 2.6, 10.8, 4.6,
        ["Propriete", "Description", "Valeur / Detail"],
        [
            ["Format entree",       "Notation moleculaire SMILES",
             "Simplified Molecular Input Line Entry System"],
            ["Labels",              "Binaires par cible",
             "0 = non-toxique  |  1 = toxique  |  NaN = non-teste"],
            ["Desequilibre",        "Classes tres majoritairement negatives",
             "Gere par class_weight='balanced' dans RF"],
            ["Val. manquantes",     "Lignes supprimees (dropna) par cible",
             "15 % (SR-MMP) a 40 % (NR-Aromatase)"],
            ["Split train/test",    "Division aleatoire fixe",
             "80 % train (6 411 mol.) / 20 % test (2 459 mol.)"],
            ["Source",              "NIH / EPA — public",
             "https://tripod.nih.gov/tox21/challenge/"],
        ],
        col_ws=[3.0, 3.8, 4.0]
    )
    add_txt(slide, 2.2, 7.15, 10.8, 0.3,
            "TOX21 — Machine Learning  |  Master Big Data S2 — 2025/2026",
            size=8, italic=True, color=GRAY_M)


def slide_06_targets(prs):
    """Slide 6 — 12 Cibles Biologiques"""
    slide = slide_base(prs)
    left_sidebar(slide, "03", "CIBLES")
    top_bar(slide)
    slide_title(slide, "Les 12 Cibles Biologiques TOX21",
                "7 recepteurs nucleaires (NR) + 5 voies de stress (SR)")

    # NR column
    add_rect(slide, 2.2, 1.15, 5.4, 0.45, NAVY, radius=True)
    add_txt(slide, 2.2, 1.17, 5.4, 0.42, "  RECEPTEURS NUCLEAIRES (NR) — 7 cibles",
            size=11, bold=True, color=WHITE)
    nr_targets = [
        ("NR-AR",         "Recepteur aux androgenes (hormone masculine)"),
        ("NR-AR-LBD",     "Domaine de liaison — recepteur androgenes"),
        ("NR-AhR",        "Recepteur AhR — dioxines / hydrocarbures"),
        ("NR-Aromatase",  "Inhibition enzyme aromatase (estrogenes)"),
        ("NR-ER",         "Recepteur aux estrogenes (perturbateur)"),
        ("NR-ER-LBD",     "Domaine de liaison — recepteur estrogenes"),
        ("NR-PPAR-gamma", "Metabolisme lipidique / diabete type 2"),
    ]
    for i, (tgt, desc) in enumerate(nr_targets):
        by = 1.65 + i * 0.75
        add_rect(slide, 2.2, by, 5.4, 0.65, WHITE, border=GRAY_L, radius=True)
        add_rect(slide, 2.2, by, 1.4, 0.65, BLUE_L, radius=False)
        add_txt(slide, 2.22, by+0.12, 1.36, 0.42, tgt,
                size=10, bold=True, color=BLUE2)
        add_txt(slide, 3.7, by+0.12, 3.85, 0.42, desc,
                size=9.5, color=GRAY_D)

    # SR column
    add_rect(slide, 7.85, 1.15, 5.25, 0.45, TEAL, radius=True)
    add_txt(slide, 7.85, 1.17, 5.25, 0.42, "  VOIES DE STRESS (SR) — 5 cibles",
            size=11, bold=True, color=WHITE)
    sr_targets = [
        ("SR-ARE",   "Reponse au stress oxydatif cellulaire"),
        ("SR-ATAD5", "Dommages ADN — genotoxicite (ATAD5)"),
        ("SR-HSE",   "Stress thermique — proteines chaperons"),
        ("SR-MMP",   "Dysfonction membrane mitochondriale"),
        ("SR-p53",   "Voie p53 — gardien du genome / cancer"),
    ]
    for i, (tgt, desc) in enumerate(sr_targets):
        by = 1.65 + i * 0.75
        add_rect(slide, 7.85, by, 5.25, 0.65, WHITE, border=GRAY_L, radius=True)
        add_rect(slide, 7.85, by, 1.25, 0.65, RGBColor(0xcc, 0xfb, 0xf1), radius=False)
        add_txt(slide, 7.87, by+0.12, 1.21, 0.42, tgt,
                size=10, bold=True, color=TEAL)
        add_txt(slide, 9.2, by+0.12, 3.85, 0.42, desc,
                size=9.5, color=GRAY_D)

    add_txt(slide, 2.2, 7.15, 10.8, 0.3,
            "TOX21 — Machine Learning  |  Master Big Data S2 — 2025/2026",
            size=8, italic=True, color=GRAY_M)


def slide_07_pipeline(prs):
    """Slide 7 — Pipeline Methodologique"""
    slide = slide_base(prs)
    left_sidebar(slide, "04", "PIPELINE")
    top_bar(slide)
    slide_title(slide, "Pipeline Methodologique",
                "De la molecule SMILES brute au deploiement de l'interface")

    steps = [
        ("📂", "Chargement\nDonnees",      "tox21.csv\n8 014 SMILES", NAVY),
        ("⚗️", "Feature\nEngineering",     "RDKit\n2 063 features",   BLUE),
        ("✂️", "Split\nTrain/Test",        "80 % / 20 %\nfixe",      BLUE2),
        ("🌲", "Entrainement\nRF x12",     "MultiOutput\nClassifier", TEAL),
        ("📊", "Evaluation\nMetriques",    "AUC-ROC\nHamming / F1",  GREEN),
        ("🤖", "SHAP +\nDeploi.",          "Streamlit\nInterface",   NAVY),
    ]
    for i, (icon, title, desc, clr) in enumerate(steps):
        bx = 2.25 + i * 1.86
        # Box
        add_rect(slide, bx, 1.3, 1.72, 2.8, clr, radius=True)
        add_rect(slide, bx, 1.3, 1.72, 2.8, clr, border=RGBColor(0x60,0xa5,0xfa), radius=True)
        add_txt(slide, bx, 1.4, 1.72, 0.65, icon,
                size=22, align=PP_ALIGN.CENTER, color=WHITE)
        add_txt(slide, bx, 2.05, 1.72, 0.65, title,
                size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_txt(slide, bx, 2.7, 1.72, 0.9, desc,
                size=9, color=BLUE_M, align=PP_ALIGN.CENTER)
        # Arrow (sauf dernier)
        if i < len(steps) - 1:
            add_txt(slide, bx+1.72, 2.45, 0.14, 0.4, "▶",
                    size=14, bold=True, color=BLUE_M, align=PP_ALIGN.CENTER)

    # Detail etapes en bas
    details = [
        ("Gestion desequilibre","class_weight='balanced' dans chaque RF",AMBER_L,AMBER),
        ("Reproductibilite","random_state=42 — resultats stables",BLUE_L,BLUE),
        ("Parallelisation","n_jobs=-1 — tous les coeurs CPU utilises",GREEN_L,GREEN),
        ("Seuil decision","Probabilite >= 0.50 = TOXIQUE",RED_L,RED),
    ]
    for i, (title, desc, bg, clr) in enumerate(details):
        bx = 2.25 + i * 2.78
        add_rect(slide, bx, 4.4, 2.6, 1.6, bg, border=clr, radius=True)
        add_txt(slide, bx+0.1, 4.5, 2.4, 0.42, title,
                size=10, bold=True, color=clr)
        add_txt(slide, bx+0.1, 4.92, 2.4, 0.9, desc,
                size=9.5, color=GRAY_D)

    add_txt(slide, 2.2, 7.15, 10.8, 0.3,
            "TOX21 — Machine Learning  |  Master Big Data S2 — 2025/2026",
            size=8, italic=True, color=GRAY_M)


def slide_08_features(prs):
    """Slide 8 — Feature Engineering Overview"""
    slide = slide_base(prs)
    left_sidebar(slide, "05", "FEATURES")
    top_bar(slide)
    slide_title(slide, "Feature Engineering",
                "Chaque molecule SMILES → vecteur numerique de 2 063 dimensions")

    # Schema 15 + 2048 = 2063
    for i, (val, lbl, clr, desc) in enumerate([
        ("15",    "Descripteurs\nPhysicochimiques", BLUE,
         "Proprietes calculees par RDKit :\nPoids mol., LogP, TPSA, LogS...\nDirectement interpretables"),
        ("+",     "",                               GRAY_M, ""),
        ("2 048", "Empreintes Morgan\nECFP4",       TEAL,
         "Vecteur binaire encodant\nla presence de sous-structures\ncirculaires (r=2 liaisons)"),
        ("=",     "",                               GRAY_M, ""),
        ("2 063", "Features Totales",               GREEN,
         "Vecteur final utilise\npour l'entrainement\ndu Random Forest"),
    ]):
        bx = 2.2 + i * 2.25
        if val in ("+","="):
            add_txt(slide, bx, 2.0, 0.55, 1.5, val,
                    size=32, bold=True, color=clr, align=PP_ALIGN.CENTER)
        else:
            add_rect(slide, bx, 1.2, 2.1, 2.6, WHITE, border=clr, radius=True)
            add_rect(slide, bx, 1.2, 2.1, 0.08, clr)
            add_txt(slide, bx, 1.35, 2.1, 0.85, val,
                    size=32, bold=True, color=clr, align=PP_ALIGN.CENTER)
            add_txt(slide, bx, 2.2, 2.1, 0.65, lbl,
                    size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
            add_txt(slide, bx+0.1, 2.9, 1.9, 1.1, desc,
                    size=9, color=GRAY_M, align=PP_ALIGN.CENTER)

    # Code snippet
    add_rect(slide, 2.2, 4.15, 10.8, 2.45, RGBColor(0x0f,0x17,0x2a), radius=True)
    code_lines = [
        ("# Calcul des 15 descripteurs physicochimiques", 9, False, GRAY_M),
        ("desc = compute_descriptors(mol)  # array de 15 floats", 9.5, False, BLUE_M),
        ("", 5, False, WHITE),
        ("# Empreintes Morgan ECFP4 — 2048 bits", 9, False, GRAY_M),
        ("fp = AllChem.GetMorganFingerprintAsBitVect(mol, radius=2, nBits=2048)", 9.5, False, BLUE_M),
        ("", 5, False, WHITE),
        ("# Concatenation → vecteur final", 9, False, GRAY_M),
        ("X = np.hstack([desc, fp])  # shape: (1, 2063)", 9.5, True, RGBColor(0x86,0xef,0xac)),
    ]
    tb = slide.shapes.add_textbox(Inches(2.45), Inches(4.3), Inches(10.3), Inches(2.2))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for (t, sz, bld, clr) in code_lines:
        if first: p = tf.paragraphs[0]; first = False
        else:     p = tf.add_paragraph()
        run = p.add_run(); run.text = t
        run.font.size = Pt(sz); run.font.bold = bld
        run.font.color.rgb = clr
        run.font.name = "Courier New"

    add_txt(slide, 2.2, 7.15, 10.8, 0.3,
            "TOX21 — Machine Learning  |  Master Big Data S2 — 2025/2026",
            size=8, italic=True, color=GRAY_M)


def slide_09_descriptors(prs):
    """Slide 9 — Descripteurs Physicochimiques"""
    slide = slide_base(prs)
    left_sidebar(slide, "05", "DESCR.")
    top_bar(slide)
    slide_title(slide, "Descripteurs Physicochimiques (15 features)",
                "Calcules avec RDKit — directement interpretables par un chimiste")

    add_table(slide, 2.2, 1.1, 10.8, 5.9,
        ["Descripteur", "Signification", "Pertinence Toxicologique"],
        [
            ["MolWt",          "Poids moleculaire (g/mol)",       "Accumulation tissulaire si > 450 g/mol"],
            ["LogP",           "Coeff. partage eau/octanol",      "Penetration membranaire et lipophilicite"],
            ["TPSA",           "Surface polaire topologique (A2)","Absorption intestinale — passage BBB"],
            ["HBondDonors",    "Donneurs de liaisons H",          "Interactions avec recepteurs proteiques"],
            ["HBondAcceptors", "Accepteurs de liaisons H",        "Affinite de fixation moleculaire"],
            ["RotatableBonds", "Liaisons rotatives",              "Flexibilite conformationnelle"],
            ["AromaticRings",  "Cycles aromatiques",             "Intercalation ADN si >= 3 cycles"],
            ["FractionCSP3",   "Fraction carbones sp3",          "Caractere plan de la molecule"],
            ["NumHeteroatoms", "Heteroatomes (N, O, S...)",      "Sites reactifs et reactivite chimique"],
            ["MolMR",          "Refractivite molaire",           "Volume moleculaire et polarisabilite"],
        ],
        col_ws=[3.0, 3.8, 4.0]
    )
    add_txt(slide, 2.2, 7.15, 10.8, 0.3,
            "TOX21 — Machine Learning  |  Master Big Data S2 — 2025/2026",
            size=8, italic=True, color=GRAY_M)


def slide_10_morgan(prs):
    """Slide 10 — Empreintes Morgan"""
    slide = slide_base(prs)
    left_sidebar(slide, "05", "MORGAN")
    top_bar(slide)
    slide_title(slide, "Empreintes Moleculaires Morgan ECFP4",
                "Encodage de la structure chimique — 2 048 bits binaires")

    # Schema Morgan
    add_rect(slide, 2.2, 1.15, 4.5, 5.5, WHITE, border=GRAY_L, radius=True)
    add_rect(slide, 2.2, 1.15, 4.5, 0.48, NAVY, radius=False)
    add_txt(slide, 2.2, 1.18, 4.5, 0.43, "  Comment ca fonctionne ?",
            size=12, bold=True, color=WHITE)
    steps_m = [
        ("1.", "Partir de chaque atome de la molecule"),
        ("2.", "Regarder le voisinage dans un rayon r=2 liaisons"),
        ("3.", "Encoder chaque sous-structure unique en 1 bit"),
        ("4.", "Produire un vecteur binaire de 2 048 positions"),
        ("",   ""),
        ("✓",  "Capture les toxicophores (groupes chimiques toxiques)"),
        ("✓",  "Standard ECFP4 utilise dans toute la litterature"),
        ("✓",  "Complementaire aux descripteurs globaux"),
    ]
    tb = slide.shapes.add_textbox(Inches(2.4), Inches(1.75), Inches(4.1), Inches(4.5))
    tf = tb.text_frame; tf.word_wrap = True; first = True
    for (num, txt_s) in steps_m:
        if first: p = tf.paragraphs[0]; first = False
        else:     p = tf.add_paragraph()
        p.space_before = Pt(5)
        if num == "": continue
        run1 = p.add_run(); run1.text = num + "  "
        run1.font.size = Pt(11); run1.font.bold = True
        run1.font.color.rgb = BLUE if num not in ("✓",) else GREEN
        run2 = p.add_run(); run2.text = txt_s
        run2.font.size = Pt(11); run2.font.color.rgb = GRAY_D

    # Vecteur visuel
    add_rect(slide, 7.0, 1.15, 6.1, 5.5, WHITE, border=GRAY_L, radius=True)
    add_rect(slide, 7.0, 1.15, 6.1, 0.48, TEAL, radius=False)
    add_txt(slide, 7.0, 1.18, 6.1, 0.43, "  Visualisation du vecteur binaire",
            size=12, bold=True, color=WHITE)
    # Grid de bits
    add_txt(slide, 7.1, 1.75, 5.9, 0.4,
            "Exemple : 2048 bits — 1 = sous-structure presente, 0 = absente",
            size=9.5, color=GRAY_M)
    # Simulation visuelle bits
    import random; random.seed(42)
    for row in range(8):
        for col in range(32):
            bit = random.random() > 0.82
            bx  = 7.1 + col * 0.185
            by  = 2.25 + row * 0.35
            clr_b = BLUE if bit else GRAY_L
            add_rect(slide, bx, by, 0.16, 0.28, clr_b, radius=False)
    add_txt(slide, 7.1, 5.15, 5.9, 0.35,
            "[ 0, 0, 1, 0, 1, 1, 0, 0, 1, ... , 0, 1, 0 ]",
            size=9.5, color=BLUE, align=PP_ALIGN.CENTER)
    add_txt(slide, 7.1, 5.5, 5.9, 0.35,
            "←————————— 2 048 bits ——————————→",
            size=9.5, color=GRAY_M, align=PP_ALIGN.CENTER)

    add_txt(slide, 2.2, 7.15, 10.8, 0.3,
            "TOX21 — Machine Learning  |  Master Big Data S2 — 2025/2026",
            size=8, italic=True, color=GRAY_M)


def slide_11_model(prs):
    """Slide 11 — Architecture du Modele"""
    slide = slide_base(prs)
    left_sidebar(slide, "06", "MODELE")
    top_bar(slide)
    slide_title(slide, "Architecture du Modele — Random Forest",
                "MultiOutputClassifier : 12 classifieurs independants en parallele")

    # Schema architecture
    # Input
    add_rect(slide, 2.3, 1.2, 2.8, 1.0, NAVY, radius=True)
    add_txt(slide, 2.3, 1.3, 2.8, 0.5, "Vecteur X", size=12, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)
    add_txt(slide, 2.3, 1.8, 2.8, 0.35, "2 063 features", size=9.5,
            color=BLUE_M, align=PP_ALIGN.CENTER)

    # Arrow
    add_txt(slide, 5.25, 1.55, 0.5, 0.5, "▶", size=22, bold=True, color=BLUE)

    # MultiOutputClassifier
    add_rect(slide, 5.85, 0.9, 3.2, 2.6, BLUE, radius=True)
    add_txt(slide, 5.85, 1.0, 3.2, 0.5, "MultiOutput\nClassifier",
            size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for j, lbl in enumerate(["RF₁","RF₂","...","RF₁₂"]):
        by = 1.55 + j * 0.48
        add_rect(slide, 6.05, by, 2.8, 0.4, BLUE2, radius=True)
        add_txt(slide, 6.05, by+0.06, 2.8, 0.3, f"RandomForest — {lbl}",
                size=9.5, color=WHITE, align=PP_ALIGN.CENTER)

    # Arrow
    add_txt(slide, 9.2, 1.55, 0.5, 0.5, "▶", size=22, bold=True, color=BLUE)

    # Outputs
    add_rect(slide, 9.85, 0.9, 2.8, 2.6, GRAY_LL, border=GRAY_L, radius=True)
    add_txt(slide, 9.85, 1.0, 2.8, 0.45, "Sorties (12 proba.)",
            size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    targets_out = [
        ("P(NR-AR)=0.72",  RED),   ("P(NR-AhR)=0.15", GREEN),
        ("P(SR-MMP)=0.88", RED),   ("P(SR-ARE)=0.31",  GREEN),
        ("...",            GRAY_M), ("P(SR-p53)=0.05",  GREEN),
    ]
    for j, (t, clr_t) in enumerate(targets_out):
        by = 1.5 + j * 0.36
        add_txt(slide, 9.95, by, 2.6, 0.32, t, size=9, bold=False, color=clr_t)

    # Comparaison modeles
    add_rect(slide, 2.3, 3.75, 10.7, 0.45, NAVY, radius=True)
    add_txt(slide, 2.3, 3.78, 10.7, 0.4, "  Comparaison des Algorithmes Evalues",
            size=12, bold=True, color=WHITE)
    add_table(slide, 2.3, 4.25, 10.7, 2.65,
        ["Algorithme", "AUC-ROC", "F1 Macro", "Hamming", "Avantage Principal"],
        [
            ["Random Forest ✓", "0.783", "0.237", "0.031", "Robuste + SHAP natif — SELECTIONNE"],
            ["Logistic Regression", "0.741", "0.198", "0.038", "Simple, rapide, moins performant"],
            ["XGBoost",         "0.798", "0.251", "0.029", "Meilleur AUC mais plus lent"],
            ["DNN (reference)",  "0.815", "0.289", "0.025", "Plus performant — moins interpretable"],
        ],
        col_ws=[2.8, 1.4, 1.4, 1.4, 3.7]
    )
    add_txt(slide, 2.2, 7.15, 10.8, 0.3,
            "TOX21 — Machine Learning  |  Master Big Data S2 — 2025/2026",
            size=8, italic=True, color=GRAY_M)


def slide_12_hyper(prs):
    """Slide 12 — Hyperparametres"""
    slide = slide_base(prs)
    left_sidebar(slide, "07", "TRAIN")
    top_bar(slide)
    slide_title(slide, "Hyperparametres et Entrainement",
                "Configuration optimale du Random Forest multi-label")

    # Table hyperparametres
    add_table(slide, 2.2, 1.15, 10.8, 3.5,
        ["Hyperparametre", "Valeur", "Justification"],
        [
            ["n_estimators",      "200 arbres",           "Equilibre performance vs. memoire/vitesse"],
            ["class_weight",      "'balanced'",           "Compensation automatique du desequilibre des classes"],
            ["random_state",      "42",                   "Garantit la reproductibilite des resultats"],
            ["max_features",      "'sqrt' (defaut)",      "Bonne generalisation — reduit la correlation entre arbres"],
            ["n_jobs",            "-1",                   "Utilise tous les coeurs CPU disponibles"],
            ["Seuil decision",    "0.50",                 "Probabilite >= 0.50 → classe 1 (TOXIQUE)"],
            ["Wrapper",           "MultiOutputClassifier","1 RF independant par cible — Binary Relevance"],
        ],
        col_ws=[3.5, 2.5, 4.8]
    )

    # Code
    add_rect(slide, 2.2, 4.85, 10.8, 2.05, RGBColor(0x0f,0x17,0x2a), radius=True)
    code_txt = [
        ("from sklearn.ensemble import RandomForestClassifier",           10, False, BLUE_M),
        ("from sklearn.multioutput import MultiOutputClassifier",          10, False, BLUE_M),
        ("",                                                               5, False, WHITE),
        ("rf = RandomForestClassifier(",                                  10, False, WHITE),
        ("    n_estimators=200, class_weight='balanced',",                10, False, WHITE),
        ("    random_state=42, n_jobs=-1",                                10, False, WHITE),
        (")",                                                              10, False, WHITE),
        ("model = MultiOutputClassifier(rf, n_jobs=-1)",                  10, False, WHITE),
        ("model.fit(X_train, Y_train)  # 12 RF entraines en parallele",  10, True, RGBColor(0x86,0xef,0xac)),
    ]
    tb = slide.shapes.add_textbox(Inches(2.45), Inches(4.98), Inches(10.3), Inches(1.8))
    tf = tb.text_frame; tf.word_wrap = True; first = True
    for (t, sz, bld, clr) in code_txt:
        if first: p = tf.paragraphs[0]; first = False
        else:     p = tf.add_paragraph()
        run = p.add_run(); run.text = t
        run.font.size = Pt(sz); run.font.bold = bld
        run.font.color.rgb = clr; run.font.name = "Courier New"

    add_txt(slide, 2.2, 7.15, 10.8, 0.3,
            "TOX21 — Machine Learning  |  Master Big Data S2 — 2025/2026",
            size=8, italic=True, color=GRAY_M)


def slide_13_results(prs):
    """Slide 13 — Resultats Globaux"""
    slide = slide_base(prs)
    left_sidebar(slide, "08", "RESULTS")
    top_bar(slide)
    slide_title(slide, "Resultats Globaux du Modele",
                "Evaluation sur 2 459 molecules de test (20 % du dataset)")

    # 3 grandes cartes metriques
    metrics = [
        ("0.783", "AUC-ROC Macro",   "Moyenne sur les 12 cibles", BLUE,  BLUE_L),
        ("0.031", "Hamming Loss",    "3.1 % de labels errones",   GREEN, GREEN_L),
        ("0.237", "F1-Score Macro",  "Limite due au desequilibre", GRAY_M, GRAY_LL),
    ]
    for i, (val, lbl, sub, clr, bg) in enumerate(metrics):
        bx = 2.3 + i * 3.7
        add_rect(slide, bx, 1.2, 3.4, 2.5, bg, border=clr, radius=True)
        add_rect(slide, bx, 1.2, 3.4, 0.08, clr)
        add_txt(slide, bx, 1.4, 3.4, 1.1, val,
                size=46, bold=True, color=clr, align=PP_ALIGN.CENTER)
        add_txt(slide, bx, 2.5, 3.4, 0.45, lbl,
                size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_txt(slide, bx, 2.95, 3.4, 0.55, sub,
                size=10, color=GRAY_M, align=PP_ALIGN.CENTER, italic=True)

    # Explication
    add_rect(slide, 2.3, 3.9, 10.8, 0.08, BLUE)
    expl = [
        (BLUE,  "AUC-ROC 0.783",  "Le modele distingue correctement toxique/non-toxique dans 78.3 % des cas"),
        (GREEN, "Hamming 0.031",  "Seulement 3.1 % des labels individuels sont incorrectement predits"),
        (AMBER, "F1 faible 0.237","Normal : le F1 est penalise par le fort desequilibre (peu de positifs)"),
    ]
    for i, (clr, titre, desc) in enumerate(expl):
        bx = 2.3 + i * 3.65
        add_rect(slide, bx, 4.1, 3.4, 2.85, WHITE, border=clr, radius=True)
        add_rect(slide, bx, 4.1, 3.4, 0.07, clr)
        add_txt(slide, bx+0.1, 4.25, 3.2, 0.5, titre,
                size=11, bold=True, color=clr)
        add_txt(slide, bx+0.1, 4.75, 3.2, 1.65, desc,
                size=10.5, color=GRAY_D)

    add_txt(slide, 2.2, 7.15, 10.8, 0.3,
            "TOX21 — Machine Learning  |  Master Big Data S2 — 2025/2026",
            size=8, italic=True, color=GRAY_M)


def slide_14_per_target(prs):
    """Slide 14 — Performance par cible"""
    slide = slide_base(prs)
    left_sidebar(slide, "08", "PERF.")
    top_bar(slide)
    slide_title(slide, "AUC-ROC par Cible Biologique",
                "Performance individuelle sur les 12 cibles TOX21")

    add_table(slide, 2.2, 1.1, 10.8, 5.8,
        ["Cible", "AUC-ROC", "Niveau", "Nb. Positifs (test)", "Interpretation"],
        [
            ["NR-AhR",       "0.861","⭐ Excellent","~180","Forte correlation structure-activite"],
            ["SR-MMP",       "0.847","⭐ Excellent","~210","Signal mitochondrial clair"],
            ["SR-ARE",       "0.826","✅ Bon",      "~195","Stress oxydatif bien predit"],
            ["SR-ATAD5",     "0.811","✅ Bon",      "~160","Dommages ADN detectables"],
            ["NR-ER",        "0.793","✅ Bon",      "~175","Perturbation estrogenique"],
            ["SR-p53",       "0.781","✅ Bon",      "~165","Voie p53 structuree"],
            ["SR-HSE",       "0.769","⚠️ Modere",  "~140","Stress thermique complexe"],
            ["NR-AR",        "0.752","⚠️ Modere",  "~120","Peu d'exemples positifs"],
            ["NR-ER-LBD",    "0.741","⚠️ Modere",  "~110","Domaine LBD specifique"],
            ["NR-AR-LBD",    "0.718","⚠️ Modere",  "~105","Tres peu d'exemples"],
            ["NR-Aromatase", "0.704","⚠️ Modere",  "~95", "Enzyme tres specifique"],
            ["NR-PPAR-gamma","0.703","⚠️ Modere",  "~90", "Signal faible"],
        ],
        col_ws=[2.5, 1.5, 1.7, 2.3, 2.8]
    )
    add_txt(slide, 2.2, 7.15, 10.8, 0.3,
            "TOX21 — Machine Learning  |  Master Big Data S2 — 2025/2026",
            size=8, italic=True, color=GRAY_M)


def slide_15_shap(prs):
    """Slide 15 — SHAP Principe"""
    slide = slide_base(prs)
    left_sidebar(slide, "09", "SHAP")
    top_bar(slide)
    slide_title(slide, "Explainabilite — SHAP",
                "SHapley Additive exPlanations — theorie des jeux cooperatifs")

    # 2 colonnes
    # Gauche : theorie
    add_rect(slide, 2.2, 1.15, 5.3, 5.6, WHITE, border=GRAY_L, radius=True)
    add_rect(slide, 2.2, 1.15, 5.3, 0.48, NAVY, radius=False)
    add_txt(slide, 2.3, 1.18, 5.1, 0.43, "Principe Theorique",
            size=13, bold=True, color=WHITE)
    lines_t = [
        "SHAP calcule la contribution de chaque",
        "feature a la prediction finale, en",
        "considerant toutes les coalitions possibles.",
        "",
        "3 proprietes fondamentales :",
        "",
        "✓ Efficacite : somme SHAP = prediction - base",
        "✓ Symetrie : features identiques → meme SHAP",
        "✓ Absence joueur fictif : SHAP=0 si no impact",
        "",
        "Implementation : SHAP TreeExplainer",
        "(optimise pour les modeles ensemblistes)",
    ]
    tb = slide.shapes.add_textbox(Inches(2.35), Inches(1.75), Inches(5.0), Inches(4.7))
    tf = tb.text_frame; tf.word_wrap = True; first = True
    for t in lines_t:
        if first: p = tf.paragraphs[0]; first = False
        else:     p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run(); run.text = t
        run.font.size = Pt(11)
        run.font.color.rgb = NAVY if t.startswith("✓") else GRAY_D
        run.font.bold = t.startswith("✓") or t.startswith("3 prop") or t.startswith("SHAP calc")

    # Droite : schema
    add_rect(slide, 7.75, 1.15, 5.35, 5.6, WHITE, border=GRAY_L, radius=True)
    add_rect(slide, 7.75, 1.15, 5.35, 0.48, BLUE, radius=False)
    add_txt(slide, 7.85, 1.18, 5.15, 0.43, "Interpretation Visuelle",
            size=13, bold=True, color=WHITE)

    # Barre SHAP simulee
    bars = [
        ("LogP=4.2",          0.82, RED),
        ("AromaticRings=3",   0.61, RED),
        ("Morgan[284]=1",     0.45, RED),
        ("Morgan[1052]=1",    0.38, RED),
        ("MolWt=387",         0.21, RED),
        ("FractionCSP3=0.11", -0.28, GREEN),
        ("HBondDonors=1",     -0.35, GREEN),
        ("TPSA=62",           -0.51, GREEN),
    ]
    base_x = 10.95
    for i, (lbl, val, clr) in enumerate(bars):
        by = 1.8 + i * 0.56
        bar_w = abs(val) * 1.8
        if val > 0:
            bx = base_x
        else:
            bx = base_x - bar_w
        add_rect(slide, bx, by+0.05, bar_w, 0.38, clr, radius=False)
        add_txt(slide, 7.85, by+0.1, 2.95, 0.35, lbl, size=8.5, color=GRAY_D)
        add_txt(slide, 12.85, by+0.1, 0.6, 0.35,
                f"{'+' if val>0 else ''}{val:.2f}", size=8.5,
                color=clr, bold=True)
    # Ligne zero
    add_rect(slide, base_x-0.02, 1.75, 0.04, len(bars)*0.56+0.15, GRAY_M)
    add_txt(slide, 10.75, 6.3, 2.5, 0.35, "Rouge=toxique | Vert=protection",
            size=9, italic=True, color=GRAY_M, align=PP_ALIGN.CENTER)

    add_txt(slide, 2.2, 7.15, 10.8, 0.3,
            "TOX21 — Machine Learning  |  Master Big Data S2 — 2025/2026",
            size=8, italic=True, color=GRAY_M)


def slide_16_shap2(prs):
    """Slide 16 — SHAP Exemple et Code"""
    slide = slide_base(prs)
    left_sidebar(slide, "09", "SHAP")
    top_bar(slide)
    slide_title(slide, "SHAP — Exemple et Implementation",
                "Explication de la toxicite du Bisphenol A sur NR-AhR")

    # Code
    add_rect(slide, 2.2, 1.15, 10.8, 2.8, RGBColor(0x0f,0x17,0x2a), radius=True)
    code_txt = [
        ("import shap", 10, False, BLUE_M),
        ("", 5, False, WHITE),
        ("# TreeExplainer sur le RF de la cible NR-AhR", 9.5, False, GRAY_M),
        ("explainer = shap.TreeExplainer(", 10, False, WHITE),
        ("    model.estimators_[2],          # RF pour NR-AhR", 10, False, WHITE),
        ("    data=background_samples,       # 80 molecules de reference", 10, False, WHITE),
        ("    feature_perturbation='interventional'", 10, False, WHITE),
        (")", 10, False, WHITE),
        ("sv = explainer.shap_values(X_bpa, check_additivity=False)", 10, False, WHITE),
        ("# sv : vecteur de 2063 valeurs SHAP pour le Bisphenol A", 9.5, True, RGBColor(0x86,0xef,0xac)),
    ]
    tb = slide.shapes.add_textbox(Inches(2.45), Inches(1.3), Inches(10.3), Inches(2.5))
    tf = tb.text_frame; tf.word_wrap = True; first = True
    for (t, sz, bld, clr) in code_txt:
        if first: p = tf.paragraphs[0]; first = False
        else:     p = tf.add_paragraph()
        run = p.add_run(); run.text = t
        run.font.size = Pt(sz); run.font.bold = bld
        run.font.color.rgb = clr; run.font.name = "Courier New"

    # Table interpretation
    add_rect(slide, 2.2, 4.1, 10.8, 0.45, NAVY, radius=True)
    add_txt(slide, 2.3, 4.13, 10.7, 0.4,
            "  Interpretation des Valeurs SHAP — Bisphenol A (NR-AhR = TOXIQUE, proba = 0.87)",
            size=12, bold=True, color=WHITE)
    add_table(slide, 2.2, 4.6, 10.8, 2.55,
        ["Feature", "Valeur", "SHAP", "Interpretation"],
        [
            ["AromaticRings",  "2",       "+0.61", "2 cycles aromatiques → forte affinite AhR"],
            ["LogP",           "3.32",    "+0.48", "Lipophile — traversee des membranes facilitee"],
            ["Morgan[2048]",   "1 (bit)", "+0.38", "Sous-structure phenolique caracteristique"],
            ["FractionCSP3",   "0.08",    "+0.31", "Molecule plane — interaction AhR favorable"],
            ["TPSA",           "40 A2",   "-0.22", "Surface polaire modeste — facteur mineur"],
        ],
        col_ws=[2.8, 1.5, 1.5, 5.0]
    )
    add_txt(slide, 2.2, 7.15, 10.8, 0.3,
            "TOX21 — Machine Learning  |  Master Big Data S2 — 2025/2026",
            size=8, italic=True, color=GRAY_M)


def slide_17_interface(prs):
    """Slide 17 — Interface Streamlit"""
    slide = slide_base(prs)
    left_sidebar(slide, "10", "UI")
    top_bar(slide)
    slide_title(slide, "Interface Utilisateur — Streamlit",
                "Application web professionnelle deployee localement")

    # Architecture
    modules = [
        ("📥", "Saisie SMILES",      "Champ texte + 4 exemples rapides\n(Bisphenol A, Dioxine, Aspirine, Ethanol)", NAVY),
        ("🔬", "Analyse ML",         "Prediction instantanee des 12 scores\nde toxicite via le modele charge", BLUE),
        ("📊", "Resultats",          "Tableau avec barres de probabilite\ncodage couleur rouge/vert par cible", BLUE2),
        ("🧬", "Structure 3D",       "Viewer 3D interactif (py3Dmol)\nRotation, zoom, styles atomiques", TEAL),
        ("🤖", "Chatbot SHAP",       "Explication en langage naturel\ndes mecanismes chimiques detectes", GREEN),
        ("📋", "Historique",         "15 dernieres molecules analysees\nBouton Recharger dans la sidebar", AMBER),
    ]
    for i, (icon, title, desc, clr) in enumerate(modules):
        col = i % 3; row = i // 3
        bx  = 2.3 + col * 3.72
        by  = 1.2 + row * 2.65
        add_rect(slide, bx, by, 3.5, 2.45, WHITE, border=clr, radius=True)
        add_rect(slide, bx, by, 3.5, 0.07, clr)
        add_txt(slide, bx, by+0.15, 3.5, 0.65, icon,
                size=26, align=PP_ALIGN.CENTER, color=clr)
        add_txt(slide, bx+0.1, by+0.85, 3.3, 0.45, title,
                size=13, bold=True, color=NAVY)
        add_txt(slide, bx+0.1, by+1.3, 3.3, 1.0, desc,
                size=9.5, color=GRAY_M)

    add_txt(slide, 2.2, 7.15, 10.8, 0.3,
            "TOX21 — Machine Learning  |  Master Big Data S2 — 2025/2026",
            size=8, italic=True, color=GRAY_M)


def slide_18_features_ui(prs):
    """Slide 18 — Fonctionnalites avancees"""
    slide = slide_base(prs)
    left_sidebar(slide, "10", "UI")
    top_bar(slide)
    slide_title(slide, "Fonctionnalites Avancees de l'Interface",
                "Dark Mode, Historique, 3D interactif, Chatbot SHAP")

    add_table(slide, 2.2, 1.15, 10.8, 3.5,
        ["Fonctionnalite", "Description", "Technologie"],
        [
            ["Prediction multi-label",   "12 scores + barres de probabilite colorees",   "Scikit-learn RF"],
            ["Chatbot SHAP",             "Explication chimique auto en langage naturel", "SHAP TreeExplainer"],
            ["Structure 3D interactive", "Rotation, zoom, styles (Boules, Batons...)",   "py3Dmol + 3Dmol.js"],
            ["Dark Mode",               "Bascule clair/sombre instantanee (sidebar)",    "CSS injection dynamique"],
            ["Historique",              "15 molecules avec bouton Recharger",            "st.session_state"],
            ["Graphique SHAP",          "Top-10 features les plus influentes",            "Matplotlib (Agg)"],
        ],
        col_ws=[3.8, 4.5, 2.5]
    )

    # Stack technique
    add_rect(slide, 2.2, 4.85, 10.8, 0.45, NAVY, radius=True)
    add_txt(slide, 2.3, 4.88, 10.7, 0.4,
            "  Stack Technique Complet",
            size=12, bold=True, color=WHITE)
    stack = [
        ("Python 3.11", BLUE),  ("Scikit-learn", BLUE), ("RDKit", TEAL),
        ("SHAP", GREEN),        ("Streamlit", RED),      ("py3Dmol", NAVY),
        ("Matplotlib", AMBER),  ("NumPy", BLUE2),        ("Pandas", TEAL),
    ]
    for i, (tech, clr) in enumerate(stack):
        bx = 2.35 + i * 1.18
        add_rect(slide, bx, 5.42, 1.08, 0.42, clr, radius=True)
        add_txt(slide, bx, 5.45, 1.08, 0.36, tech,
                size=8.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Commande lancement
    add_rect(slide, 2.2, 6.0, 10.8, 0.85, RGBColor(0x0f,0x17,0x2a), radius=True)
    add_txt(slide, 2.45, 6.1, 10.3, 0.65,
            "streamlit run app.py    →    http://localhost:8501",
            size=14, bold=True, color=RGBColor(0x86,0xef,0xac),
            align=PP_ALIGN.CENTER)

    add_txt(slide, 2.2, 7.15, 10.8, 0.3,
            "TOX21 — Machine Learning  |  Master Big Data S2 — 2025/2026",
            size=8, italic=True, color=GRAY_M)


def slide_19_conclusion(prs):
    """Slide 19 — Conclusion et Perspectives"""
    slide = slide_base(prs)
    left_sidebar(slide, "11", "CONCL.")
    top_bar(slide)
    slide_title(slide, "Conclusion et Perspectives",
                "Bilan du projet et axes d'amelioration futurs")

    # Bilan gauche
    add_rect(slide, 2.2, 1.15, 5.3, 3.5, WHITE, border=GRAY_L, radius=True)
    add_rect(slide, 2.2, 1.15, 5.3, 0.48, GREEN, radius=False)
    add_txt(slide, 2.3, 1.18, 5.1, 0.43, "✅ Ce Que Nous Avons Realise",
            size=12, bold=True, color=WHITE)
    bilan = [
        "Pipeline ML complet de bout en bout",
        "AUC-ROC macro 0.783 (competitif litterature)",
        "Feature engineering chimique : 2 063 dims",
        "Explainabilite SHAP par cible biologique",
        "Interface Streamlit professionnelle deployee",
        "Visualisation 3D moleculaire interactive",
        "Dark Mode + Historique des analyses",
    ]
    for i, b in enumerate(bilan):
        add_txt(slide, 2.3, 1.75+i*0.35, 5.1, 0.35,
                f"▶  {b}", size=10.5, color=GRAY_D)

    # Perspectives droite
    add_rect(slide, 7.75, 1.15, 5.35, 3.5, WHITE, border=GRAY_L, radius=True)
    add_rect(slide, 7.75, 1.15, 5.35, 0.48, BLUE, radius=False)
    add_txt(slide, 7.85, 1.18, 5.15, 0.43, "🔭 Perspectives Futures",
            size=12, bold=True, color=WHITE)
    persp = [
        ("GNN — Graph Neural Networks",  "Exploiter le graphe moleculaire"),
        ("XGBoost / LightGBM",           "Ameliorer AUC sur cibles rares"),
        ("Authentification Google OAuth", "Acces personnalise par user"),
        ("Deploiement Streamlit Cloud",   "Accessibilite publique"),
        ("SMOTE sur-echantillonnage",     "Ameliorer le F1-Score"),
        ("Validation externe",           "Test sur d'autres datasets"),
    ]
    for i, (title, desc) in enumerate(persp):
        by = 1.75 + i * 0.41
        add_txt(slide, 7.85, by, 5.15, 0.22, f"▷  {title}",
                size=10.5, bold=True, color=BLUE)
        add_txt(slide, 8.1, by+0.22, 4.85, 0.2, desc,
                size=9, color=GRAY_M)

    # Chiffres cles bas
    add_rect(slide, 2.2, 4.85, 10.8, 0.06, BLUE)
    key_nums = [("0.783","AUC-ROC"),("0.031","Hamming"),("2063","Features"),
                ("12","Cibles"),("8014","Molecules"),("200","Arbres RF")]
    for i, (v, l) in enumerate(key_nums):
        bx = 2.3 + i * 1.83
        add_rect(slide, bx, 5.05, 1.65, 1.85, NAVY, radius=True)
        add_txt(slide, bx, 5.15, 1.65, 0.85, v,
                size=22, bold=True, color=BLUE_M, align=PP_ALIGN.CENTER)
        add_txt(slide, bx, 5.95, 1.65, 0.35, l,
                size=9, color=GRAY_M, align=PP_ALIGN.CENTER)

    add_txt(slide, 2.2, 7.15, 10.8, 0.3,
            "TOX21 — Machine Learning  |  Master Big Data S2 — 2025/2026",
            size=8, italic=True, color=GRAY_M)


def slide_20_merci(prs):
    """Slide 20 — Merci + References"""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Fond
    add_rect(slide, 0, 0, 13.33, 7.5, NAVY)
    add_rect(slide, 0, 0, 13.33, 7.5, RGBColor(0x0a,0x14,0x24))
    # Deco cercles
    add_rect(slide, 0, 5.5, 4.5, 4.5, BLUE2, radius=False)
    c1 = slide.shapes[-1]; c1.fill.fore_color.rgb = RGBColor(0x0d,0x1e,0x40)
    c1.line.color.rgb = RGBColor(0x0d,0x1e,0x40)
    add_rect(slide, 10.5, -0.5, 4.0, 4.0, BLUE2, radius=False)
    c2 = slide.shapes[-1]; c2.fill.fore_color.rgb = RGBColor(0x0d,0x1e,0x40)
    c2.line.color.rgb = RGBColor(0x0d,0x1e,0x40)

    # MERCI
    add_txt(slide, 1.5, 0.6, 10.33, 1.5,
            "MERCI",
            size=68, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, 3.5, 2.1, 6.33, 0.06, BLUE)
    add_txt(slide, 1.5, 2.25, 10.33, 0.55,
            "Nous remercions toutes les personnes qui ont contribue a ce projet",
            size=13, italic=True, color=BLUE_M, align=PP_ALIGN.CENTER)

    # Auteur + superviseur
    add_rect(slide, 3.2, 2.95, 3.0, 1.2, RGBColor(0x0d,0x1b,0x35), radius=True)
    add_rect(slide, 3.2, 2.95, 3.0, 1.2, RGBColor(0x0d,0x1b,0x35), border=BLUE, radius=True)
    add_txt(slide, 3.2, 3.05, 3.0, 0.38, "Prepare par :", size=9, color=GRAY_M, align=PP_ALIGN.CENTER)
    add_txt(slide, 3.2, 3.43, 3.0, 0.5, "[Prenom NOM]", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_txt(slide, 3.2, 3.93, 3.0, 0.2, "Master Big Data S2", size=9, color=GRAY_M, align=PP_ALIGN.CENTER)

    add_rect(slide, 7.13, 2.95, 3.0, 1.2, RGBColor(0x0d,0x1b,0x35), radius=True)
    add_rect(slide, 7.13, 2.95, 3.0, 1.2, RGBColor(0x0d,0x1b,0x35), border=BLUE, radius=True)
    add_txt(slide, 7.13, 3.05, 3.0, 0.38, "Encadre par :", size=9, color=GRAY_M, align=PP_ALIGN.CENTER)
    add_txt(slide, 7.13, 3.43, 3.0, 0.5, "Pr. [Superviseur]", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_txt(slide, 7.13, 3.93, 3.0, 0.2, "Universite Hassan II", size=9, color=GRAY_M, align=PP_ALIGN.CENTER)

    # References cles
    add_rect(slide, 1.5, 4.35, 10.33, 0.4, BLUE, radius=False)
    add_txt(slide, 1.5, 4.37, 10.33, 0.38, "  References Principales",
            size=11, bold=True, color=WHITE)
    refs = [
        "[1] TOX21 Challenge — NIH/EPA. https://tripod.nih.gov/tox21/challenge/",
        "[2] Huang R. et al. Tox21Challenge — Frontiers in Environmental Science, 2016.",
        "[3] Lundberg S., Lee S. A Unified Approach to Interpreting Model Predictions. NeurIPS, 2017.",
        "[4] Pedregosa F. et al. Scikit-learn: Machine Learning in Python. JMLR, 2011.",
        "[5] Landrum G. et al. RDKit: Open-source cheminformatics. https://www.rdkit.org",
    ]
    for i, r in enumerate(refs):
        add_txt(slide, 1.65, 4.85+i*0.42, 10.05, 0.38, r, size=9.5, color=GRAY_M)

    # Bas
    add_rect(slide, 0, 6.85, 13.33, 0.65, BLUE)
    add_txt(slide, 0, 6.9, 13.33, 0.55,
            "Universite Hassan II  —  Faculte des Sciences Ben M'Sick  —  Master Big Data S2  —  2025/2026",
            size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    print("Generation des 20 slides...")
    slide_01_cover(prs)       ; print("  [1/20] Cover")
    slide_02_toc(prs)         ; print("  [2/20] Sommaire")
    slide_03_context(prs)     ; print("  [3/20] Contexte")
    slide_04_multilabel(prs)  ; print("  [4/20] Multi-Label")
    slide_05_dataset(prs)     ; print("  [5/20] Dataset")
    slide_06_targets(prs)     ; print("  [6/20] 12 Cibles")
    slide_07_pipeline(prs)    ; print("  [7/20] Pipeline")
    slide_08_features(prs)    ; print("  [8/20] Features Overview")
    slide_09_descriptors(prs) ; print("  [9/20] Descripteurs")
    slide_10_morgan(prs)      ; print(" [10/20] Morgan")
    slide_11_model(prs)       ; print(" [11/20] Architecture")
    slide_12_hyper(prs)       ; print(" [12/20] Hyperparametres")
    slide_13_results(prs)     ; print(" [13/20] Resultats Globaux")
    slide_14_per_target(prs)  ; print(" [14/20] Perf. par Cible")
    slide_15_shap(prs)        ; print(" [15/20] SHAP Principe")
    slide_16_shap2(prs)       ; print(" [16/20] SHAP Exemple")
    slide_17_interface(prs)   ; print(" [17/20] Interface")
    slide_18_features_ui(prs) ; print(" [18/20] Fonctionnalites UI")
    slide_19_conclusion(prs)  ; print(" [19/20] Conclusion")
    slide_20_merci(prs)       ; print(" [20/20] Merci")

    out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                       "Presentation_TOX21_ML.pptx")
    prs.save(out)
    print(f"\n[OK] PPT genere : {out}")


if __name__ == "__main__":
    main()
